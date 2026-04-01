import os
import re
import zipfile
import io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from core.models import Song, FontCalculation
from config import settings


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
VERSE_LEFT = Inches(0)
VERSE_TOP = Inches(0)
VERSE_WIDTH = SLIDE_WIDTH
VERSE_HEIGHT = SLIDE_HEIGHT
BG_COLOR = RGBColor(0, 0, 0)
TEXT_COLOR = RGBColor(255, 255, 255)


class PPTGeneratorError(Exception):
    pass


def _embed_thumbnail_manually(pptx_path: str, song: Song, font_size: int):
    try:
        order = song.get_effective_order()
        if not order:
            return
        
        first_key = order[0]
        first_text = song.get_text_for_key(first_key)
        if not first_text:
            return
        
        img = _generate_thumbnail_image(first_text, font_size)
        
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='JPEG', quality=85)
        img_buffer.seek(0)
        img_data = img_buffer.getvalue()
        
        with zipfile.ZipFile(pptx_path, 'r') as zip_read:
            files = zip_read.namelist()
            file_contents = {}
            for file in files:
                if file != 'docProps/thumbnail.jpeg':
                    file_contents[file] = zip_read.read(file)
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for file, content in file_contents.items():
                zip_write.writestr(file, content)
            zip_write.writestr('docProps/thumbnail.jpeg', img_data)
        
    except Exception as e:
        print(f"Could not embed thumbnail: {e}")


def _generate_thumbnail_image(text: str, font_size: int) -> Image.Image:
    width, height = 320, 180
    img = Image.new('RGB', (width, height), color=(0, 0, 0))
    draw = ImageDraw.Draw(img)
    
    thumbnail_font_size = 20
    try:
        font = ImageFont.truetype("calibri.ttf", thumbnail_font_size)
    except:
        try:
            font = ImageFont.truetype("arial.ttf", thumbnail_font_size)
        except:
            font = ImageFont.load_default()
    
    lines = text.split("\n")
    max_line_length = max(len(line) for line in lines) if lines else 0
    
    adjusted_font_size = thumbnail_font_size
    if max_line_length > 30:
        adjusted_font_size = int(thumbnail_font_size * (30 / max_line_length))
        adjusted_font_size = max(adjusted_font_size, 12)
        
        try:
            font = ImageFont.truetype("calibri.ttf", adjusted_font_size)
        except:
            try:
                font = ImageFont.truetype("arial.ttf", adjusted_font_size)
            except:
                font = ImageFont.load_default()
    
    line_height = adjusted_font_size * 1.3
    max_lines = int(height / line_height) - 1
    lines = lines[:max_lines]
    
    total_text_height = len(lines) * line_height
    start_y = (height - total_text_height) // 2
    
    for i, line in enumerate(lines):
        bbox = draw.textbbox((0, 0), line, font=font)
        text_width = bbox[2] - bbox[0]
        x = (width - text_width) // 2
        y = start_y + i * line_height
        draw.text((x, y), line, fill=(255, 255, 255), font=font)
    
    return img


def generate_pptx(song: Song, font_calc: FontCalculation,
                  output_path: str, use_caps: bool = False) -> str:
    try:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]
        
        config_refrain = settings.get_section("refrain")
        start_marker = config_refrain.get("start_marker", "/:")
        end_marker = config_refrain.get("end_marker", ":/")
        
        order = song.get_effective_order()

        # Calculează indexul ultimului slide cu text (pentru Amin corect)
        last_idx_with_text = -1
        for i, key in enumerate(order):
            if song.get_text_for_key(key):
                last_idx_with_text = i

        for i, key in enumerate(order):
            text = song.get_text_for_key(key)
            if not text:
                continue

            if use_caps:
                text = _apply_caps(text, start_marker, end_marker)

            is_refrain = key in song.refrains
            if is_refrain:
                # Doar normalizeaza spatierea daca marcajele exista deja
                # Nu mai adauga marcaje automate - userul le editeaza manual
                lines = text.split('\n')

                for line_idx in range(len(lines)):
                    lines[line_idx] = lines[line_idx].replace(': /', ':/')

                text = '\n'.join(lines)

            is_last = (i == last_idx_with_text)
            _add_verse_slide(prs, blank_layout, text, font_calc.size_pt,
                           add_amin=is_last, use_caps=use_caps)
        
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        _embed_thumbnail_manually(output_path, song, font_calc.size_pt)
        
        return os.path.abspath(output_path)
    
    except Exception as e:
        raise PPTGeneratorError(f"Eroare la generarea PPT: {str(e)}")


def _add_verse_slide(prs, layout, text: str, font_size: int, 
                    add_amin: bool = False, use_caps: bool = False):
    slide = prs.slides.add_slide(layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    txBox = slide.shapes.add_textbox(VERSE_LEFT, VERSE_TOP, 
                                      VERSE_WIDTH, VERSE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.color.rgb = TEXT_COLOR
            run.font.bold = False
        
        p.space_after = Pt(font_size * 0.2)
        p.space_before = Pt(0)
    
    if add_amin:
        _add_amin_to_slide(slide, use_caps)


def _add_amin_to_slide(slide, use_caps: bool = False):
    amin_config = settings.get_section("amin_slide")
    amin_text = "AMIN" if use_caps else "Amin"
    
    right_cm = 1
    bottom_cm = 1
    text_box_width = Inches(3)
    text_box_height = Inches(1)
    
    left = SLIDE_WIDTH - Cm(right_cm) - text_box_width
    top = SLIDE_HEIGHT - Cm(bottom_cm) - text_box_height
    
    txBox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = tf.paragraphs[0]
    p.text = amin_text
    p.font.name = amin_config.get("font_family", "Calibri")
    p.font.size = Pt(amin_config.get("font_size", 22))
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = amin_config.get("is_bold", False)
    p.font.italic = amin_config.get("is_italic", False)
    p.alignment = PP_ALIGN.RIGHT


def _apply_caps(text: str, start_marker: str, end_marker: str) -> str:
    text = text.replace(start_marker, "<<<START>>>")
    text = text.replace(end_marker, "<<<END>>>")
    text = text.upper()
    text = text.replace("<<<START>>>", start_marker)
    text = text.replace("<<<END>>>", end_marker)
    return text


def generate_filename(title: str) -> str:
    name = title.split(' - ')[0] if ' - ' in title else title
    name = name.replace('-', ' ')
    
    replacements = {
        'ă': 'a', 'â': 'a', 'î': 'i', 'ș': 's', 'ț': 't', 'ş': 's', 'ţ': 't',
        'Ă': 'A', 'Â': 'A', 'Î': 'I', 'Ș': 'S', 'Ț': 'T', 'Ş': 'S', 'Ţ': 'T'
    }
    for old, new in replacements.items():
        name = name.replace(old, new)
    
    name = re.sub(r'[^a-zA-Z\s]', '', name)
    name = re.sub(r'\s+', ' ', name.strip())
    
    words = name.split()
    name = ' '.join(words[:6])
    
    return f"{name}.pptx"
